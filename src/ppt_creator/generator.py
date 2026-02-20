from pptx import Presentation

class PPTGenerator:
    def __init__(self, title="My Presentation"):
        self.prs = Presentation()
        self.title = title

    def add_title_slide(self, title_text, subtitle_text=""):
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text

    def add_content_slide(self, title_text, bullet_points):
        slide_layout = self.prs.slide_layouts[1]  # Bullet point layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    def save(self, filename):
        self.prs.save(filename)
        print(f"Presentation saved as {filename}")
